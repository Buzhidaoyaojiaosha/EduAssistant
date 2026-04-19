import os
import json
import asyncio
from typing import List, Dict
from PyPDF2 import PdfReader
from flask import current_app
from pptx import Presentation
from app.ext import knowledge_base_collection, rag_chunk_collection
from langchain.text_splitter import RecursiveCharacterTextSplitter
from app.ext import embedding_fn
from app.react.tools_register import register_as_tool
from docx import Document
from moviepy import VideoFileClip
import dashscope
from dashscope import MultiModalConversation



class RAGService:
    @staticmethod
    def extract_text_from_file(file_path: str) -> str:
        """从文件中提取文本内容。支持 PDF、PPTX、DOCX 格式。视频文件请走 Qwen 多模态路径。"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件未找到: {file_path}")

        if file_path.endswith('.pdf'):
            return RAGService._extract_text_from_pdf(file_path)
        elif file_path.endswith('.pptx'):
            return RAGService._extract_text_from_pptx(file_path)
        elif file_path.endswith('.docx'):
            return RAGService._extract_text_from_docx(file_path)
        else:
            raise ValueError("不支持的文件格式")

    @staticmethod
    def _extract_text_from_pdf(file_path: str) -> str:
        """从 PDF 文件中提取文本。"""
        text = ''
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + '\n'
        return text

    @staticmethod
    def _extract_text_from_pptx(file_path: str) -> str:
        """从 PPTX 文件中提取文本。"""
        text = ''
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        return text
    @staticmethod
    def _extract_text_from_docx(file_path: str) -> str:
        """从 DOCX 文件中提取文本。"""
        doc = Document(file_path)
        text = ''
        for para in doc.paragraphs:
            text += para.text + '\n'
        
        print("以下是文件提取的文本内容:")
        print(text)
        return text
    

    @staticmethod
    def _extract_knowledge_from_video_via_qwen(file_path: str, title: str):
        """使用 DashScope 千问多模态API直接分析视频，提取结构化知识步骤。

        Args:
            file_path: 视频文件本地路径
            title: 视频标题

        Returns:
            tuple: (plain_text: str, word_doc: Document) 纯文本用于向量库存储，Word文档用于用户下载
        """
        dashscope.api_key = os.environ.get("DASHSCOPE_API_KEY")
        if not dashscope.api_key:
            raise ValueError("DASHSCOPE_API_KEY 未配置")

        abs_path = os.path.abspath(file_path)
        local_file_url = f"file://{abs_path}"

        prompt_text = (
            f"你是一个教学视频分析专家。请仔细观看这段名为《{title}》的教学视频。\n"
            "请将视频中的教学内容拆解为连续的、结构化的知识步骤。\n"
            "每个步骤包含：步骤编号(step_number)、步骤名称(action_name)、详细说明(details)。\n"
            "警告：绝对不要在 details 中生成或编造任何图片链接、Markdown图片语法。只输出纯文本说明。\n"
            "要求严格按照JSON数组格式返回，不带```json标记。\n"
            '格式示例：[{"step_number": 1, "action_name": "动作名称", "details": "详细说明"}]'
        )

        messages = [
            {
                "role": "user",
                "content": [
                    {"video": local_file_url},
                    {"text": prompt_text}
                ]
            }
        ]

        response = asyncio.run(
            asyncio.to_thread(
                MultiModalConversation.call,
                model='qwen-vl-plus',
                messages=messages
            )
        )

        if response.status_code != 200:
            raise ValueError(f"千问API调用失败: {response.message}")

        content = response.output.choices[0].message.content[0].get('text', '')
        content = content.strip()
        # 清理可能的 Markdown 代码块标记
        if content.startswith("```json"):
            content = content[7:]
        if content.startswith("```"):
            content = content[3:]
        if content.endswith("```"):
            content = content[:-3]
        content = content.strip()

        steps = json.loads(content)

        # 构建纯文本（用于 ChromaDB 向量库存储）
        plain_text_parts = []
        for step in steps:
            plain_text_parts.append(
                f"步骤{step['step_number']}: {step['action_name']}\n{step['details']}"
            )
        plain_text = "\n\n".join(plain_text_parts)

        # 构建 Word 文档（用于用户下载）
        doc = Document()
        doc.add_heading(f"《{title}》视频分析报告", level=1)
        doc.add_paragraph(f"本文档由智能教学系统根据视频自动分析生成。")
        doc.add_paragraph("")
        for step in steps:
            doc.add_heading(f"步骤{step['step_number']}: {step['action_name']}", level=2)
            doc.add_paragraph(step['details'])
            doc.add_paragraph("")

        return plain_text, doc

    @staticmethod
    def split_text(text: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
        """将文本切分为多个片段。"""
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ".", "。", "！", "？", "，", " "]
        )
        return text_splitter.split_text(text)

    @staticmethod
    def process_and_store_file(file_path: str, file_url: str, title: str, course_id: int = None,
                               category: str = None, tags: List[str] = None):
        """
        处理文件，提取文本，切分内容，并存储到向量数据库。

        Args:
            file_path: 文件本地路径
            file_url: 文件oss路径
            title: 知识条目标题
            course_id: 关联的课程 ID
            category: 分类
            tags: 标签列表

        Returns:
            str or None: 视频文件返回生成的Word文档路径，其他文件返回None
        """
        word_file_path = None

        # 视频文件走千问多模态路径
        if file_path.endswith(('.mp4', '.avi', '.mov', '.mkv')):
            text, word_doc = RAGService._extract_knowledge_from_video_via_qwen(file_path, title)
            # 保存 Word 文档
            word_dir = os.path.join(current_app.config['UPLOAD_FOLDER'], 'generated_docs')
            os.makedirs(word_dir, exist_ok=True)
            safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
            word_filename = f"{safe_title}_analysis.docx"
            word_file_path = os.path.join(word_dir, word_filename)
            word_doc.save(word_file_path)
        else:
            text = RAGService.extract_text_from_file(file_path)
        print(text)
        if not isinstance(text, str):
            raise ValueError(f"提取的文本不是字符串类型，文件路径: {file_path}")
        chunks = RAGService.split_text(text)

        # 保存切分片段到 rag_chunk_collection
        ids = []
        documents = []
        metadatas = []
        for index, chunk in enumerate(chunks):
            doc_id = f"{title}_片段{index}"
            ids.append(doc_id)
            documents.append(chunk)
            metadata = {
                "title": f"{title}_片段{index}",
                "type": "1",
                "course_id": course_id,
                "category": category,
                "tags": ",".join(tags) if tags else "",
                "is_chunk": True,
                "chunk_index": index,
                "source_file": file_url
            }
            metadatas.append(metadata)

        # 批量添加到向量数据库（chromadb 根据 collection 的 embedding_function 自动生成 embeddings）
        if documents:
            rag_chunk_collection.add(
                ids=ids,
                documents=documents,
                metadatas=metadatas
            )

        return word_file_path

        # # 保存切分片段
        # for index, chunk in enumerate(chunks):
        #     KnowledgeBaseService.add_knowledge(
        #         title=f"{title}_片段{index}",
        #         type=1,
        #         content=chunk,
        #         course_id=course_id,
        #         category=category,
        #         tags=tags,
        #         is_chunk=True,
        #         chunk_index=index,
        #         source_file=file_url
        #     )

    @staticmethod
    def delete_chunks_by_source_url(source_url):
        """
        根据 source_url 删除向量数据库中 is_chunk 为 true 的记录
        :param source_url: 文件的源 URL
        """
        try:
            # 查询满足条件的记录的 ID
            results = knowledge_base_collection.get(
                where={"$and": [{"is_chunk": True}, {"source_file": source_url}]}
            )
            ids = results.get("ids", [])

            if ids:
                # 删除匹配的记录
                knowledge_base_collection.delete(ids=ids)
                rag_chunk_collection.delete(ids=ids)
                return True
            return False
        except Exception as e:
            current_app.logger.error(f"删除向量数据库记录失败: {str(e)}")
            return False


    @register_as_tool(roles=["student", "teacher"])
    @staticmethod
    def search_docs(query, top_k=3):
        """基于给定查询在向量数据库中搜索相关文档片段。
        Args:
            query: 用于搜索的查询文本。
            top_k: 返回的相关文档片段的最大数量，默认为3。

        Returns:
            str: 包含匹配文档片段及其来源信息的字符串，各片段间用换行符分隔。

        Raises:
            Exception: 当向量数据库查询过程中出现错误时抛出。
        """

        embedding = embedding_fn([query])  # embedding_fn 是 OpenAI 或 sentence-transformers

        results = rag_chunk_collection.query(
            query_embeddings=embedding,
            n_results=top_k,
            where={"is_chunk": True}  # 只搜索切片数据
        )

        # 提取文本片段和来源
        docs = []
        for doc, meta in zip(results["documents"][0], results["metadatas"][0]):
            docs.append(f"《{meta['title']}》第{meta['chunk_index']}段：{doc}")

        return "\n".join(docs)
