from typing import Dict, List
from datetime import datetime
import os
import tempfile
from werkzeug.utils import secure_filename
import io
import base64
import subprocess
import json
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from flask import current_app, jsonify
from app.react.tools_register import register_as_tool
from app.utils.logging import logger
from app.models.course import Course
from app.services.knowledge_base_service import KnowledgeBaseService
from app.services.course_service import CourseService
from app.utils.llm.deepseek import chat_deepseek
import re
from pptx import Presentation
from reportlab.lib import colors
from reportlab.platypus import ListFlowable, HRFlowable

# 配置背景图片路径
BACKGROUND_IMAGE_PATH = r"app\static\images\ppt_background.png"

def generate_teaching_material(course_id: int, material_type: str, selected_knowledge_ids: List[int] = None) -> Dict:
    """生成AI驱动的备课资料
    
    Args:
        course_id (int): 课程ID
        material_type(str):资料类型（syllabus/ppt）
        selected_knowledge_ids (List[int], optional): 选中的知识库ID列表
        
    Returns:
        Dict: 包含生成的备课资料内容和文件信息
    """
    try:
        course = Course.get_by_id(course_id)
        if not course:
            return {"error": "课程不存在"}
        
        # 1. 获取本课程的知识库资料
        relevant_docs = _get_selected_knowledge_content(course_id, selected_knowledge_ids)
        
        # 2. 构造DeepSeek API提示词
        if material_type == "syllabus":
            prompt = f"""
你是一位经验丰富的教师，需要为《{course.name}》课程制作详细的备课提纲。

相关课程资料：
{relevant_docs}

请根据以上课程资料，生成一份完整的备课提纲，包含以下内容：

1. 课程概述（课程性质、地位和作用）
2. 教学目标（知识目标、能力目标、素质目标）
3. 教学重点难点
4. 教学内容安排（详细列出主要知识点和讲解要点）
5. 实训练习与指导（设计具体的练习题和指导方案）
6. 教学方法与策略
7. 教学资源需求（教材、设备、软件等）
8. 课堂组织形式
9. 评估方式与标准
10. 课后拓展与作业

要求：
- 内容要具体详实，可直接用于教学
- 充分结合课程资料中的知识点和内容
- 教学安排要合理，循序渐进
- 实训练习要有针对性和实用性
- 语言专业规范，格式清晰
- 体现理论与实践相结合的教学理念

请用中文输出，结构化呈现。
"""
        else:  # PPT
            prompt = f"""
你是一位经验丰富的教师，需要为《{course.name}》课程制作教学PPT的Markdown内容。

相关课程资料：
{relevant_docs}

请按照Marp格式生成教学PPT的Markdown内容，要求如下：

1. 使用Marp语法，以---分隔幻灯片
2. 包含12-15张幻灯片
3. 每张幻灯片内容简洁明了，要点清晰
4. 内容结构：
   - 封面页：课程标题和基本信息
   - 课程介绍
   - 教学目标
   - 重点难点
   - 主要内容（按知识点分页，每页3-5个要点）
   - 实践练习
   - 教学方法
   - 评估标准
   - 总结与作业

格式要求：
- 每个要点使用-或*开头的列表
- 重要内容用**粗体**标出
- 适当使用# ## ###标题层级
- 每张幻灯片内容控制在合理范围内

请严格按照以下Marp Markdown格式输出：

---
marp: true
theme: default
class: invert
paginate: true
---

# 课程标题
## 副标题信息

---

# 第一页标题
- 要点1
- 要点2
- 要点3

---

...（后续页面按此格式）

请确保输出的是标准的Marp Markdown格式，可以直接被marp-cli处理。
"""
        
        messages = [
            {"role": "system", "content": "你是一位专业的教育专家和课程设计师，擅长制作教学材料。"},
            {"role": "user", "content": prompt}
        ]
        
        # 3. 调用DeepSeek API生成内容
        ai_content = chat_deepseek(messages)
        if not ai_content:
            raise ValueError("DeepSeek API返回空响应")
        
        # 4. 生成文件
        if material_type == "syllabus":
            file_base64, filename = _generate_pdf_content(
                title=f"《{course.name}》备课提纲",
                content=ai_content,
                course=course
            )
            file_type = "pdf"
        else:
            # 使用Marp生成PPT
            file_base64, filename = _generate_marp_ppt(
                title=f"《{course.name}》教学PPT",
                markdown_content=ai_content,
                course=course
            )
            file_type = "pptx"
        
        return {
            "content": ai_content,
            "file_base64": file_base64,
            "filename": filename,
            "title": f"《{course.name}》备课材料",
            "download_ready": True,
            "material_type": material_type,
            "file_type": file_type
        }
        
    except Exception as e:
        logger.error(f"生成教学材料失败: {str(e)}")
        return {"error": f"生成失败: {str(e)}"}

def _get_selected_knowledge_content(course_id: int, selected_ids: List[int] = None) -> str:
    """获取选中的知识库内容作为AI生成参考
    
    Args:
        course_id (int): 课程ID
        selected_ids (List[int]): 选中的知识库ID列表
        
    Returns:
        str: 格式化的知识库内容
    """
    try:
        # 如果没有提供选中的ID，获取所有知识库内容
        if not selected_ids:
            all_knowledge = CourseService.get_knowledge_base_by_course(course_id=course_id)
        else:
            # 获取选中的知识库内容
            all_knowledge = KnowledgeBaseService.get_knowledge_by_ids(selected_ids)
        
        if not all_knowledge:
            return "没有选择任何课程资料，将基于基础模板生成内容。"
        
        # 格式化内容
        formatted_content = []
 
        for item in all_knowledge:
            # 判断内容类型：下载链接还是纯文本
            if item.content.startswith(('http://', 'https://', 'ftp://')):
                # 这是一个下载链接
                content_desc = f"文件链接: {item.content}"
                content_type = "可下载资源"
            else:
                # 这是纯文本内容
                content_preview = item.content[:500] + "..." if len(item.content) > 500 else item.content
                content_desc = f"文本内容: {content_preview}"
                content_type = "纯文本资源"
            
            formatted_item = f"""
【{item.title}】
类别: {item.category or '未分类'}
资源类型: {content_type}
内容: {content_desc}
创建时间: {item.created_at.strftime('%Y-%m-%d')}
---
"""
            formatted_content.append(formatted_item)
        
        result = f"以下是《课程ID:{course_id}》的相关教学资料（共{len(all_knowledge)}项）：\n\n"
        result += "\n".join(formatted_content)
        
        return result
        
    except Exception as e:
        logger.error(f"获取选中知识库内容失败: {str(e)}")
        return "获取课程资料失败，将基于基础模板生成内容。"
    
def _remove_slides(prs, remove_first=5, remove_last=1):
    """从PPT中删除指定数量的开头和结尾幻灯片"""
    try:
        total_slides = len(prs.slides)
        if total_slides <= remove_first + remove_last:
            logger.warning(f"PPT总页数{total_slides}不足{remove_first+remove_last}页，跳过删除操作")
            return prs
        
        # 删除前N页
        for _ in range(remove_first):
            if prs.slides:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                if slides:
                    xml_slides.remove(slides[0])
        
        # 删除最后M页
        for _ in range(remove_last):
            if prs.slides:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                if slides:
                    xml_slides.remove(slides[-1])
        
        return prs
    except Exception as e:
        logger.error(f"删除PPT幻灯片失败: {str(e)}")
        return prs
    
def _generate_marp_ppt(title: str, markdown_content: str, course: Course) -> tuple:
    """使用Marp生成美观的PPT文件并返回base64编码
    
    Args:
        title (str): PPT标题
        markdown_content (str): Marp格式的Markdown内容
        course (Course): 课程对象
        
    Returns:
        tuple: (base64_content, filename)
    """
    try:
        # 检查是否安装了marp-cli
        if not _check_marp_installed():
            logger.warning("marp-cli未安装，回退到简单PPT生成")
            return _generate_simple_ppt_content(title, markdown_content, course)
        
        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            # 生成时间戳和文件名
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            md_filename = f"slides_{timestamp}.md"
            pptx_filename = f"{title}_{timestamp}.pptx"
            
            md_path = os.path.join(temp_dir, md_filename)
            pptx_path = os.path.join(temp_dir, pptx_filename)
            
            # 处理Markdown内容，确保符合Marp格式
            processed_content = _process_marp_markdown(markdown_content, title, course)
            
            # 写入Markdown文件
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(processed_content)
            
           # 使用marp-cli转换为PPTX
            try:
                import platform
                is_windows = platform.system().lower() == 'windows'
            
                # 使用正确的包名 @marp-team/marp-cli
                cmd = [
                    'npx', '@marp-team/marp-cli',  
                    '--allow-local-files',
                    '--pptx',
                    '--output', pptx_path,
                    md_path
                ]
            
                # 设置环境变量
                env = os.environ.copy()
                env['NODE_NO_WARNINGS'] = '1'
            
                # Windows 下使用 shell=True
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=60,
                    cwd=temp_dir,
                    env=env,
                    shell=is_windows
                )
                
                if result.returncode != 0:
                    logger.error(f"Marp转换失败: {result.stderr}")
                    logger.error(f"标准输出: {result.stdout}")
                    raise subprocess.CalledProcessError(result.returncode, cmd, result.stderr)
                
                # 检查输出文件是否存在
                if not os.path.exists(pptx_path):
                    logger.error(f"输出文件不存在: {pptx_path}")
                    logger.error(f"临时目录内容: {os.listdir(temp_dir)}")
                    raise FileNotFoundError("Marp转换完成但输出文件不存在")
                
                # 读取生成的PPTX文件
                if Presentation is not None:
                    try:
                        # 打开PPT并删除指定页
                        prs = Presentation(pptx_path)
                        prs = _remove_slides(prs, remove_first=5, remove_last=1)
                    
                        # 保存修改后的PPT
                        modified_pptx_path = os.path.join(temp_dir, f"modified_{pptx_filename}")
                        prs.save(modified_pptx_path)
                        pptx_path = modified_pptx_path  # 使用修改后的文件路径
                    except Exception as e:
                        logger.error(f"删除PPT页面失败，使用原始文件: {str(e)}")
                # 读取生成的PPTX文件
                with open(pptx_path, 'rb') as f:
                    pptx_content = f.read()
                
                # 转换为base64
                pptx_base64 = base64.b64encode(pptx_content).decode('utf-8')
                
                logger.info(f"成功使用Marp生成PPT: {pptx_filename}, 文件大小: {len(pptx_content)} bytes")
                return pptx_base64, pptx_filename
                
            except (subprocess.TimeoutExpired, subprocess.CalledProcessError, FileNotFoundError) as e:
                logger.error(f"Marp命令执行失败: {str(e)}")
                logger.error(f"命令: {' '.join(cmd) if isinstance(cmd, list) else cmd}")
                # 回退到简单PPT生成
                return _generate_simple_ppt_content(title, markdown_content, course)
                
    except Exception as e:
        logger.error(f"Marp PPT生成失败: {str(e)}")
        # 回退到简单PPT生成
        return _generate_simple_ppt_content(title, markdown_content, course)


def _check_marp_installed() -> bool:
    """改进的Marp安装检测函数"""
    try:
        # 设置环境变量来抑制Node.js警告
        env = os.environ.copy()
        env['NODE_NO_WARNINGS'] = '1'
        
        result = subprocess.run(
            ['npx', 'marp', '--version'],
            capture_output=True,
            text=True,
            timeout=15,  # 增加超时时间
            env=env,     # 使用修改后的环境变量
            shell=True   # 在Windows上使用shell
        )
        
        # 记录详细信息用于调试
        logger.info(f"Marp版本检查 - 返回码: {result.returncode}")
        logger.info(f"标准输出: {result.stdout.strip()}")
        if result.stderr:
            logger.info(f"错误输出: {result.stderr.strip()}")
        
        # 检查是否包含版本信息
        if result.returncode == 0 and '@marp-team/marp-cli' in result.stdout:
            logger.info("Marp CLI 检测成功")
            return True
        else:
            logger.warning(f"Marp CLI 检测失败，返回码: {result.returncode}")
            return False
            
    except subprocess.TimeoutExpired:
        logger.error("Marp 版本检查超时")
        return False
    except FileNotFoundError:
        logger.error("找不到 npx 命令")
        return False
    except Exception as e:
        logger.error(f"检查 Marp 安装时出错: {str(e)}")
        return False
    
def _process_marp_markdown(content: str, title: str, course: Course) -> str:
    """处理和优化Marp Markdown内容"""

    # 如果内容不包含Marp头部，添加默认配置
    if not content.strip().startswith('---\nmarp: true'):
        marp_header = f"""---
marp: true
theme: default
class: lead
paginate: true
backgroundColor: #fff
style: |
  section {{
    background-size: cover;
    background-position: center;
    background-repeat: no-repeat;
    color: #333;
  }}
  section.lead {{
    text-align: center;
    color: #fff;
    text-shadow: 2px 2px 4px rgba(0,0,0,0.5);
  }}
  section.lead h1 {{
    color: #fff;
    font-size: 3em;
    text-shadow: 2px 2px 4px rgba(0,0,0,0.7);
  }}
  section.lead h2 {{
    color: #f0f0f0;
    text-shadow: 1px 1px 2px rgba(0,0,0,0.5);
  }}
  section h1 {{
    color: #2c3e50;
    border-bottom: 3px solid #3498db;
    padding-bottom: 10px;
    background: rgba(255,255,255,0.9);
    padding: 15px;
    border-radius: 10px;
    margin-bottom: 20px;
  }}
  section h2 {{
    color: #2c3e50;
    background: rgba(255,255,255,0.8);
    padding: 10px;
    border-radius: 8px;
  }}
  section ul, section p {{
    background: rgba(255,255,255,0.85);
    padding: 15px;
    border-radius: 8px;
    margin-bottom: 15px;
    box-shadow: 0 2px 5px rgba(0,0,0,0.1);
  }}
  section li {{
    margin-bottom: 8px;
    line-height: 1.6;
  }}
  section strong {{
    color: #e74c3c;
    font-weight: bold;
  }}
---

"""
        # 如果内容已经有封面页，直接添加头部
        if content.strip().startswith('#'):
            content = marp_header + content
        else:
            # 添加默认封面页
            cover_slide = f"""
# {title}

## {course.name}

**教师:** {course.teacher.name if hasattr(course, 'teacher') and course.teacher else '未指定'}  
**生成时间:** {datetime.now().strftime('%Y年%m月%d日')}

---

"""
            content = marp_header + cover_slide + content
    
    # 优化内容格式
    lines = content.split('\n')
    processed_lines = []
    
    for line in lines:
        # 确保幻灯片分隔符独占一行
        if line.strip() == '---':
            processed_lines.append('')
            processed_lines.append('---')
            processed_lines.append('')
        else:
            processed_lines.append(line)
    
    return '\n'.join(processed_lines)

def _clean_text(text: str) -> str:
    """清理和格式化文本内容"""
    # 替换特殊字符
    text = text.replace('**', '')  # 移除Markdown加粗标记
    text = text.replace('*', '')   # 移除Markdown斜体标记
    text = text.replace('`', '')   # 移除代码标记
    text = text.replace('"', '"')  # 替换引号
    text = text.replace('"', '"')
    text = text.replace(''', "'")
    text = text.replace(''', "'")
    text = text.replace('…', '...')
    text = text.replace('—', '-')
    text = text.replace('–', '-')
    text = text.replace('―', '-')
    # 清理多余的空白字符
    text = ' '.join(text.split())
    return text

def _process_markdown_content(content: str) -> List[Dict]:
    """处理Markdown内容为结构化数据"""
    sections = []
    current_section = None
    current_list = []
    
    lines = content.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            if current_list:
                if current_section:
                    if 'content' not in current_section:
                        current_section['content'] = []
                    current_section['content'].extend(current_list)
                current_list = []
            continue
        
        # 处理标题
        if line.startswith('# '):
            if current_section and current_list:
                if 'content' not in current_section:
                    current_section['content'] = []
                current_section['content'].extend(current_list)
                current_list = []
            
            if current_section:
                sections.append(current_section)
            
            current_section = {
                'type': 'h1',
                'text': _clean_text(line[2:]),
                'content': []
            }
        elif line.startswith('## '):
            if current_section and current_list:
                if 'content' not in current_section:
                    current_section['content'] = []
                current_section['content'].extend(current_list)
                current_list = []
            
            if current_section:
                sections.append(current_section)
            
            current_section = {
                'type': 'h2',
                'text': _clean_text(line[3:]),
                'content': []
            }
        # 处理列表项
        elif line.startswith(('- ', '* ', '+ ')):
            current_list.append({
                'type': 'bullet',
                'text': _clean_text(line[2:])
            })
        # 处理有序列表
        elif re.match(r'^\d+\.\s', line):
            number = line.split('.')[0]
            text = line[len(number)+2:]
            current_list.append({
                'type': 'number',
                'number': number,
                'text': _clean_text(text)
            })
        # 处理普通段落
        else:
            if current_list:
                if current_section:
                    if 'content' not in current_section:
                        current_section['content'] = []
                    current_section['content'].extend(current_list)
                current_list = []
            
            if current_section:
                current_section['content'].append({
                    'type': 'paragraph',
                    'text': _clean_text(line)
                })
    
    # 处理最后的部分
    if current_list and current_section:
        if 'content' not in current_section:
            current_section['content'] = []
        current_section['content'].extend(current_list)
    
    if current_section:
        sections.append(current_section)
    
    return sections

def _generate_simple_ppt_content(title: str, content: str, course: Course) -> tuple:
    """简单PPT生成方法（原有的python-pptx实现作为后备）"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # 创建PPT对象
        prs = Presentation()
        
        # 设置默认字体（支持中文）
        # 添加封面幻灯片
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        subtitle = slide.placeholders[1]
        subtitle.text = f"课程: {course.name}\n教师: {getattr(course.teacher, 'name', '未指定') if hasattr(course, 'teacher') else '未指定'}\n生成时间: {datetime.now().strftime('%Y年%m月%d日')}"
        
        # 解析Markdown内容生成幻灯片
        slides_content = _parse_markdown_to_slides(content)
        
        for slide_data in slides_content:
            # 添加内容幻灯片
            slide_layout = prs.slide_layouts[1]  # 标题和内容布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get('title', '无标题')
            
            # 设置内容
            if 'content' in slide_data and slide_data['content']:
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.clear()
                
                for item in slide_data['content']:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(18)
        
        # 保存到临时文件并转换为base64
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            ppt_path = tmp_file.name
            prs.save(ppt_path)
        
        try:
            with open(ppt_path, 'rb') as f:
                ppt_content = f.read()
            os.unlink(ppt_path)
            
            ppt_base64 = base64.b64encode(ppt_content).decode('utf-8')
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"{title}_{timestamp}.pptx"
            
            return ppt_base64, filename
        finally:
            # 确保临时文件被删除
            if os.path.exists(ppt_path):
                os.unlink(ppt_path)
        
    except Exception as e:
        logger.error(f"简单PPT生成失败: {str(e)}")
        # 最终回退：生成文本文件
        return _generate_text_fallback(title, content, course)

def _parse_markdown_to_slides(content: str) -> List[Dict]:
    """解析Markdown内容为幻灯片数据"""
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    
    for line in lines:
        line = line.strip()
        
        if not line:
            continue
            
        # 检测幻灯片分隔符
        if line == '---':
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': '', 'content': []}
            continue
            
        # 检测标题
        if line.startswith('# '):
            if current_slide is None:
                current_slide = {'title': '', 'content': []}
            if not current_slide['title']:
                current_slide['title'] = line[2:].strip()
            else:
                current_slide['content'].append(line[2:].strip())
        elif line.startswith('## '):
            if current_slide is None:
                current_slide = {'title': '', 'content': []}
            current_slide['content'].append(line[3:].strip())
        # 检测列表项
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide is None:
                current_slide = {'title': '内容', 'content': []}
            current_slide['content'].append(line[2:].strip())
        # 普通文本
        elif line and not line.startswith('marp:') and not line.startswith('theme:') and not line.startswith('class:'):
            if current_slide is None:
                current_slide = {'title': '内容', 'content': []}
            if not current_slide['title']:
                current_slide['title'] = line
            else:
                current_slide['content'].append(line)
    
    # 添加最后一张幻灯片
    if current_slide:
        slides.append(current_slide)
    
    return slides

def _generate_text_fallback(title: str, content: str, course: Course) -> tuple:
    """最终回退方案：生成文本文件"""
    try:
        txt_content = f"{title}\n\n"
        txt_content += f"课程：{course.name}\n"
        txt_content += f"教师：{getattr(course.teacher, 'name', '未指定') if hasattr(course, 'teacher') else '未指定'}\n"
        txt_content += f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        txt_content += "=" * 50 + "\n\n"
        txt_content += content
        
        txt_base64 = base64.b64encode(txt_content.encode('utf-8')).decode('utf-8')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{title}_{timestamp}.txt"
        
        return txt_base64, filename
        
    except Exception as e:
        logger.error(f"文本回退也失败: {str(e)}")
        # 返回最基础的错误响应
        error_content = f"生成失败: {str(e)}"
        error_base64 = base64.b64encode(error_content.encode('utf-8')).decode('utf-8')
        return error_base64, "error.txt"

def _generate_pdf_content(title: str, content: str, course: Course) -> tuple:
    """生成格式化的PDF文件并返回base64编码"""
    try:
        # 创建内存缓冲区
        buffer = io.BytesIO()
        
        # 生成文件名
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{title}_{timestamp}.pdf"
        
        # 注册中文字体
        try:
            pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
            chinese_font = 'STSong-Light'
        except:
            try:
                font_path = os.path.join(os.path.dirname(__file__), 'fonts', 'simhei.ttf')
                if not os.path.exists(font_path):
                    font_path = 'C:/Windows/Fonts/simhei.ttf'
                    if not os.path.exists(font_path):
                        font_path = '/usr/share/fonts/truetype/wqy/wqy-microhei.ttc'
                
                pdfmetrics.registerFont(TTFont('ChineseFont', font_path))
                chinese_font = 'ChineseFont'
            except Exception as font_error:
                logger.error(f"字体注册失败: {str(font_error)}")
                chinese_font = 'Helvetica'
        
        # 创建PDF文档
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            topMargin=0.6*inch,
            bottomMargin=0.6*inch,
            leftMargin=0.8*inch,
            rightMargin=0.8*inch,
            encoding='utf-8'
        )
        
        # 创建样式
        styles = getSampleStyleSheet()
        
        # 标题样式
        title_style = ParagraphStyle(
            'TitleStyle',
            parent=styles['Heading1'],
            fontSize=24,
            leading=32,
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName=chinese_font,
            textColor=colors.HexColor('#2c3e50'),
            borderWidth=2,
            borderColor=colors.HexColor('#3498db'),
            borderPadding=20,
            backColor=colors.HexColor('#f8f9fa')
        )
        
        # 章节标题样式
        heading1_style = ParagraphStyle(
            'Heading1Style',
            parent=styles['Heading2'],
            fontSize=18,
            leading=24,
            spaceBefore=20,
            spaceAfter=12,
            fontName=chinese_font,
            textColor=colors.HexColor('#2c3e50'),
            borderWidth=0,
            borderColor=colors.HexColor('#3498db'),
            borderPadding=10,
            leftIndent=0,
            firstLineIndent=0,
            borderRadius=5,
            backColor=colors.HexColor('#f8f9fa')
        )
        
        # 章节副标题样式
        heading2_style = ParagraphStyle(
            'Heading2Style',
            parent=styles['Heading3'],
            fontSize=14,
            leading=20,
            spaceBefore=15,
            spaceAfter=10,
            fontName=chinese_font,
            textColor=colors.HexColor('#34495e'),
            borderWidth=0,
            borderRadius=5,
            leftIndent=20,
            firstLineIndent=0,
            backColor=colors.HexColor('#f8f9fa')
        )
        
        # 正文样式
        body_style = ParagraphStyle(
            'BodyStyle',
            parent=styles['Normal'],
            fontSize=12,
            leading=18,
            spaceBefore=6,
            spaceAfter=6,
            fontName=chinese_font,
            textColor=colors.HexColor('#2c3e50'),
            alignment=TA_LEFT,
            wordWrap='CJK',
            firstLineIndent=24,
            leftIndent=12
        )
        
        # 列表样式
        list_style = ParagraphStyle(
            'ListStyle',
            parent=body_style,
            fontSize=12,
            leading=18,
            leftIndent=40,
            bulletIndent=20,
            spaceBefore=6,
            spaceAfter=6,
            textColor=colors.HexColor('#2c3e50'),
            bulletFontName=chinese_font,
            bulletFontSize=10
        )
        
        # 课程信息样式
        info_style = ParagraphStyle(
            'InfoStyle',
            parent=body_style,
            fontSize=10,
            leading=14,
            alignment=TA_CENTER,
            textColor=colors.HexColor('#666666'),
            backColor=colors.HexColor('#f8f9fa'),
            borderPadding=10,
            borderRadius=5
        )
        
        # 处理Markdown内容
        sections = _process_markdown_content(content)
        
        # 构建文档内容
        story = []
        
        # 添加主标题
        story.append(Spacer(1, 30))
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 20))
        
        # 添加课程信息
        course_info = f"""
        <para alignment="center">
        <font color="#666666">课程名称:</font> {course.name} &nbsp;|&nbsp; 
        <font color="#666666">课程代码:</font> {course.code} &nbsp;|&nbsp; 
        <font color="#666666">生成时间:</font> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
        </para>
        """
        story.append(Paragraph(course_info, info_style))
        story.append(Spacer(1, 30))
        
        # 添加分隔线
        story.append(HRFlowable(
            width="100%",
            thickness=1,
            lineCap='round',
            color=colors.HexColor('#e0e0e0'),
            spaceBefore=10,
            spaceAfter=20
        ))
        
        # 处理每个部分
        for section in sections:
            if section['type'] == 'h1':
                story.append(Spacer(1, 20))
                story.append(Paragraph(section['text'], heading1_style))
                story.append(Spacer(1, 10))
            elif section['type'] == 'h2':
                story.append(Spacer(1, 15))
                story.append(Paragraph(section['text'], heading2_style))
                story.append(Spacer(1, 8))
            
            # 处理内容
            if 'content' in section:
                current_list_items = []
                
                for item in section['content']:
                    if item['type'] in ['bullet', 'number']:
                        current_list_items.append(
                            Paragraph(item['text'], list_style)
                        )
                    else:  # paragraph
                        if current_list_items:
                            story.append(ListFlowable(
                                current_list_items,
                                bulletType='bullet',
                                bulletColor=colors.HexColor('#3498db'),
                                bulletFontName=chinese_font,
                                bulletFontSize=8,
                                leftIndent=40,
                                bulletDedent=20,
                                spaceBefore=10,
                                spaceAfter=10
                            ))
                            current_list_items = []
                        story.append(Paragraph(item['text'], body_style))
                
                # 处理最后的列表项
                if current_list_items:
                    story.append(ListFlowable(
                        current_list_items,
                        bulletType='bullet',
                        bulletColor=colors.HexColor('#3498db'),
                        bulletFontName=chinese_font,
                        bulletFontSize=8,
                        leftIndent=40,
                        bulletDedent=20,
                        spaceBefore=10,
                        spaceAfter=10
                    ))
        
        # 添加页脚
        def add_footer(canvas, doc):
            canvas.saveState()
            canvas.setFont(chinese_font, 8)
            canvas.setFillColor(colors.HexColor('#999999'))
            footer_text = f"第 {canvas.getPageNumber()} 页"
            page_width = doc.pagesize[0]
            canvas.drawCentredString(page_width/2.0, 30, footer_text)
            canvas.restoreState()
        
        # 构建PDF
        doc.build(story, onFirstPage=add_footer, onLaterPages=add_footer)
        
        # 获取PDF内容并转换为base64
        pdf_content = buffer.getvalue()
        buffer.close()
        
        pdf_base64 = base64.b64encode(pdf_content).decode('utf-8')
        
        return pdf_base64, filename
        
    except Exception as e:
        logger.error(f"PDF生成失败: {str(e)}")
        return _generate_text_fallback(title, content, course)

# 添加保存教学资料到知识库的功能
def save_teaching_material_to_kb(course_id: int, title: str, file_base64: str, filename: str) -> Dict:
    """将教学资料保存到知识库并使用OSS存储"""
    try:
        course = Course.get_by_id(course_id)
        if not course:
            return {"error": "课程不存在"}
        
        # 解码base64
        file_data = base64.b64decode(file_base64)
        
        # 创建有意义的文件名
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        file_ext = os.path.splitext(filename)[1].lower()
        material_type = "备课PPT" if file_ext == '.pptx' else "备课大纲"
        
        # 创建安全的文件名：课程名_类型_时间戳.扩展名
        base_name = f"{course.name}_{material_type}_{timestamp}"
        safe_base_name = secure_filename(base_name)
        safe_filename = f"{safe_base_name}{file_ext}"
        
        # 创建临时文件
        with tempfile.NamedTemporaryFile(
            delete=False, 
            suffix=file_ext,
            prefix=safe_base_name + "_"  # 添加前缀
        ) as tmp_file:
            tmp_file.write(file_data)
            tmp_file_path = tmp_file.name
        
        try:
            # 上传文件到OSS
            file_url = KnowledgeBaseService.upload_file_to_oss(tmp_file_path)
            
            # 根据文件扩展名确定类型
            file_ext = os.path.splitext(filename)[1].lower()
            if file_ext == '.pdf':
                kb_type = "2"  # 2: pdf
            elif file_ext == '.pptx':
                kb_type = "3"  # 3: pptx
            else:
                kb_type = "4"  # 4: 其他
                
            # 添加知识条目到知识库
            new_knowledge = KnowledgeBaseService.add_knowledge(
                title=title,
                type=kb_type,
                content=file_url,  # 使用OSS返回的URL
                course_id=course_id,
                category="备课资料",
                tags=["文档"],
                is_chunk=False,
                chunk_index=None,
                source_file=file_url  # 源文件也是OSS URL
            )
            
            return {"success": True, "knowledge_id": new_knowledge.id}
        finally:
            # 确保删除临时文件
            if os.path.exists(tmp_file_path):
                os.unlink(tmp_file_path)
                
    except Exception as e:
        current_app.logger.error(f"保存教学资料到知识库失败: {str(e)}")
        return {"success": False, "error": str(e)}

def get_marp_themes() -> List[str]:
    """获取可用的Marp主题列表"""
    return [
        'default',
        'gaia',
        'uncover',
        'academic',
        'corporate',
        'gradient'
    ]

def create_custom_marp_theme() -> str:
    """创建自定义Marp主题CSS"""
    return """
/* @theme custom */

@import 'default';

:root {
  --color-background: #fafafa;
  --color-foreground: #333;
  --color-highlight: #0066cc;
  --color-dimmed: #666;
}

section {
  background-color: var(--color-background);
  color: var(--color-foreground);
  font-family: 'Microsoft YaHei', 'SimHei', sans-serif;
  padding: 70px;
}

section.lead h1 {
  text-align: center;
  color: var(--color-highlight);
  border-bottom: 3px solid var(--color-highlight);
  padding-bottom: 20px;
}

section h1 {
  color: var(--color-highlight);
  border-bottom: 2px solid #eee;
  padding-bottom: 10px;
}

section h2 {
  color: var(--color-highlight);
}

section ul {
  font-size: 1.1em;
  line-height: 1.6;
}

section li {
  margin-bottom: 8px;
}

section strong {
  color: var(--color-highlight);
  font-weight: bold;
}

section em {
  color: var(--color-dimmed);
  font-style: italic;
}

section code {
  background: #f5f5f5;
  border-radius: 3px;
  padding: 2px 6px;
}

footer {
  color: var(--color-dimmed);
}
"""

# 导出的主要函数（保持原有接口兼容性）
__all__ = [
    'generate_teaching_material',
    'get_marp_themes',
    'create_custom_marp_theme',
    'save_teaching_material_to_kb',
]